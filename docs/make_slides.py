"""Generate trt-lightnet technical presentation — clean white design."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Palette ──────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_DARK  = RGBColor(0x44, 0x44, 0x44)
GRAY_MID   = RGBColor(0x88, 0x88, 0x88)
GRAY_LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
GRAY_LINE  = RGBColor(0xDD, 0xDD, 0xDD)
GREEN      = RGBColor(0x00, 0x96, 0x64)   # accent
BLUE       = RGBColor(0x00, 0x6A, 0xC1)
CODE_BG    = RGBColor(0xF7, 0xF7, 0xF7)
CODE_FG    = RGBColor(0x24, 0x6A, 0x00)   # dark green for code

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Primitives ───────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=WHITE, line=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def bullet_box(slide, lines, l, t, w, h,
               size=12, color=BLACK, bg=None, line_color=None):
    if bg:
        rect(slide, l, t, w, h, fill=bg,
             line=line_color, lw=Pt(0.75) if line_color else Pt(0))
    tb = slide.shapes.add_textbox(
        Inches(l + 0.12), Inches(t + 0.1),
        Inches(w - 0.24), Inches(h - 0.2))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color


def page_header(slide, title, subtitle=None):
    """Thin green top accent + large title."""
    rect(slide, 0, 0, 13.33, 0.07, fill=GREEN)
    txt(slide, title, 0.5, 0.2, 12.0, 0.72,
        size=28, bold=True, color=BLACK)
    if subtitle:
        txt(slide, subtitle, 0.5, 0.88, 12.0, 0.3,
            size=13, color=GRAY_MID, italic=True)
    rect(slide, 0.5, 1.25, 12.33, 0.02, fill=GRAY_LINE)


def page_footer(slide, n, total=16):
    rect(slide, 0, 7.3, 13.33, 0.2, fill=GRAY_LIGHT)
    txt(slide, "TensorRT-LightNet  |  TIER IV, Inc.",
        0.4, 7.31, 9, 0.18, size=9, color=GRAY_MID)
    txt(slide, f"{n} / {total}", 12.4, 7.31, 0.8, 0.18,
        size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def section_label(slide, text, l, t, w, color=GREEN):
    rect(slide, l, t, w, 0.35, fill=color)
    txt(slide, f"  {text}", l, t, w, 0.35,
        size=12, bold=True, color=WHITE)


def table(slide, headers, rows, l, t, w, col_widths,
          hdr_color=GREEN, row_size=12):
    cx = [l]
    for cw in col_widths[:-1]:
        cx.append(cx[-1] + cw)
    # header
    for ci, (hdr, cw, x) in enumerate(zip(headers, col_widths, cx)):
        rect(slide, x, t, cw, 0.35, fill=hdr_color)
        txt(slide, hdr, x + 0.08, t + 0.02, cw - 0.12, 0.31,
            size=11, bold=True, color=WHITE)
    # rows
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else GRAY_LIGHT
        for ci, (val, cw, x) in enumerate(zip(row, col_widths, cx)):
            rect(slide, x, t + 0.35 + ri * 0.32, cw, 0.32,
                 fill=bg, line=GRAY_LINE, lw=Pt(0.5))
            bold = (ci == 0)
            color = GREEN if ci == 0 else GRAY_DARK
            txt(slide, val,
                x + 0.08, t + 0.37 + ri * 0.32,
                cw - 0.12, 0.28,
                size=row_size, bold=bold, color=color)


# ════════════════════════════════════════════════════════════════
# 1 — Title
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
rect(sl, 0, 0, 13.33, 0.07, fill=GREEN)
rect(sl, 0, 7.43, 13.33, 0.07, fill=GREEN)

txt(sl, "TensorRT-LightNet",
    1.0, 2.0, 11.33, 1.3,
    size=48, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

rect(sl, 4.2, 3.45, 4.93, 0.04, fill=GREEN)

txt(sl, "High-Efficiency Real-Time CNN Inference for Edge AI",
    1.0, 3.6, 11.33, 0.6,
    size=18, color=GRAY_DARK, align=PP_ALIGN.CENTER)

txt(sl, "Technical Overview",
    1.0, 4.3, 11.33, 0.4,
    size=14, color=GRAY_MID, italic=True, align=PP_ALIGN.CENTER)

txt(sl, "TIER IV, Inc.   ·   github.com/tier4/trt-lightnet",
    1.0, 6.8, 11.33, 0.35,
    size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 2 — Agenda
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Agenda")
page_footer(sl, 2)

items = [
    "01  What is trt-lightnet?",
    "02  Key Features",
    "03  Requirements & Installation",
    "04  Quick Start — Engine Build & Inference",
    "05  Configuration Files",
    "06  Multitask Inference",
    "07  Two-Stage Hierarchical Detection",
    "08  Uncertainty Estimation & BEV Projection",
    "09  LiDAR Range Image Segmentation",
    "10  Python API  (pylightnet)",
    "11  Performance Optimization",
    "12  Command-Line Reference",
    "13  Project Structure",
    "14  References",
]
for i, item in enumerate(items):
    col, row = i % 2, i // 2
    lx = 0.5 + col * 6.4
    ty = 1.4 + row * 0.52
    rect(sl, lx, ty, 6.0, 0.44, fill=GRAY_LIGHT, line=GRAY_LINE, lw=Pt(0.5))
    num = item[:2]
    rest = item[4:]
    rect(sl, lx, ty, 0.44, 0.44, fill=GREEN)
    txt(sl, num, lx, ty, 0.44, 0.44,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, rest, lx + 0.5, ty + 0.05, 5.4, 0.34, size=13, color=BLACK)


# ════════════════════════════════════════════════════════════════
# 3 — What is trt-lightnet?
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "What is trt-lightnet?")
page_footer(sl, 3)

txt(sl,
    "trt-lightnet is a TensorRT-based CNN inference framework optimized for NVIDIA edge AI devices "
    "(Jetson Xavier, Orin, etc.). It combines LightNet — a lightweight neural network architecture — "
    "with TensorRT to deliver real-time multitask perception.",
    0.5, 1.35, 12.33, 0.85, size=15, color=GRAY_DARK)

pillars = [
    ("LightNet",    "Lightweight CNN\narchitecture for\nedge devices"),
    ("+",           ""),
    ("TensorRT",    "NVIDIA inference\noptimizer & GPU\nacceleration"),
    ("=",           ""),
    ("trt-lightnet","Real-time multitask\nperception pipeline\nfor autonomous driving"),
]
bw, bh = 2.4, 1.6
xs = [0.5, 3.0, 3.5, 6.0, 6.5]
for i, (title, desc) in enumerate(pillars):
    lx = xs[i]
    if title in ("+", "="):
        txt(sl, title, lx, 2.35, 0.5, 0.8,
            size=28, bold=True, color=GRAY_MID, align=PP_ALIGN.CENTER)
    else:
        color = GREEN if title == "trt-lightnet" else BLUE
        rect(sl, lx, 2.2, bw, 0.42, fill=color)
        txt(sl, title, lx, 2.2, bw, 0.42,
            size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, lx, 2.62, bw, bh - 0.42, fill=GRAY_LIGHT, line=GRAY_LINE)
        txt(sl, desc, lx + 0.1, 2.68, bw - 0.2, bh - 0.55,
            size=13, color=GRAY_DARK, align=PP_ALIGN.CENTER)

txt(sl, "Supports:", 0.5, 4.1, 2.0, 0.4, size=13, bold=True, color=BLACK)
caps = [
    "Object Detection",
    "Semantic Segmentation",
    "Depth Estimation",
    "LiDAR Range Image Segmentation",
    "Two-stage Hierarchical Classification",
    "Entropy-based Uncertainty Estimation",
]
for i, cap in enumerate(caps):
    col, row = i % 3, i // 3
    lx = 0.5 + col * 4.2
    ty = 4.55 + row * 0.48
    rect(sl, lx, ty, 0.08, 0.28, fill=GREEN)
    txt(sl, cap, lx + 0.2, ty, 3.8, 0.38, size=13, color=GRAY_DARK)


# ════════════════════════════════════════════════════════════════
# 4 — Key Features
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Key Features")
page_footer(sl, 4)

features = [
    ("Multitask Inference",        "Detection + Segmentation + Depth in a single forward pass"),
    ("Multi-Precision",            "FP32 / FP16 / INT8 with per-layer precision control"),
    ("2:4 Structured Sparsity",   "50% weight reduction, maintained accuracy on Ampere+ GPUs"),
    ("NVDLA Support",              "Offload to Deep Learning Accelerator on Xavier / Orin"),
    ("Hierarchical Detection",     "Two-stage subnet for fine-grained classification (e.g. TLR color)"),
    ("Uncertainty Estimation",     "Entropy-based confidence maps from softmax outputs"),
    ("BEV Projection",             "Bird's Eye View from monocular depth maps"),
    ("Range Image Segmentation",   "LiDAR point cloud → range image → semantic segmentation"),
    ("Python Bindings (pylightnet)","ctypes-based wrapper for scripting and integration"),
]
for i, (title, desc) in enumerate(features):
    col, row = i % 3, i // 3
    lx = 0.4 + col * 4.3
    ty = 1.4 + row * 1.82
    rect(sl, lx, ty, 4.05, 1.65, fill=WHITE, line=GRAY_LINE, lw=Pt(0.75))
    rect(sl, lx, ty, 0.06, 1.65, fill=GREEN)
    txt(sl, title, lx + 0.18, ty + 0.12, 3.75, 0.4,
        size=13, bold=True, color=BLACK)
    rect(sl, lx + 0.18, ty + 0.52, 3.6, 0.02, fill=GRAY_LINE)
    txt(sl, desc, lx + 0.18, ty + 0.65, 3.75, 0.9,
        size=12, color=GRAY_DARK)


# ════════════════════════════════════════════════════════════════
# 5 — Requirements & Installation
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Requirements & Installation")
page_footer(sl, 5)

section_label(sl, "Requirements", 0.5, 1.35, 5.4)
table(sl,
      ["Dependency", "Version"],
      [
          ("CUDA",      "11.0+  (tested: 11.4 / 11.7 / 11.8 / 12.x)"),
          ("TensorRT",  "8.5 / 8.6 / 10.x  (TRT10 supported)"),
          ("CMake",     "3.10+"),
          ("GCC",       "≤ 11.x"),
          ("OS",        "Ubuntu 22.04 / Jetson JetPack 5.1"),
          ("Docker",    "24.0+  +  NVIDIA Container Toolkit 1.14+"),
      ],
      0.5, 1.7, 5.4, [1.5, 3.9])

section_label(sl, "Local Install", 6.4, 1.35, 6.5, color=BLUE)
code = [
    "# 1. Clone with submodules",
    "git clone --recurse-submodules \\",
    "  git@github.com:tier4/trt-lightnet.git",
    "",
    "# 2. Install system libraries",
    "sudo apt install -y \\",
    "  libgflags-dev libboost-all-dev \\",
    "  libopencv-dev libeigen3-dev \\",
    "  nlohmann-json3-dev libssl-dev",
    "",
    "# 3. Build & install",
    "mkdir build && cd build",
    "cmake ../ && make -j$(nproc)",
    "sudo make install",
]
bullet_box(sl, code, 6.4, 1.7, 6.5, 4.8,
           size=11, color=CODE_FG, bg=CODE_BG)

section_label(sl, "Docker", 0.5, 5.85, 5.4, color=GRAY_DARK)
docker = [
    "docker build -f Dockerfile_x86     -t trt-lightnet:latest .  # x86_64",
    "docker build -f Dockerfile_aarch64 -t trt-lightnet:latest .  # Jetson",
    "docker run -it --gpus all       trt-lightnet:latest   # x86_64",
    "docker run -it --runtime=nvidia trt-lightnet:latest   # Jetson",
]
bullet_box(sl, docker, 0.5, 6.2, 12.4, 0.98,
           size=11, color=CODE_FG, bg=CODE_BG)


# ════════════════════════════════════════════════════════════════
# 6 — Quick Start
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Quick Start", "Engine Build → Inference")
page_footer(sl, 6)

section_label(sl, "Step 1 — Build TensorRT Engine  (ONNX → .engine)", 0.5, 1.35, 12.33)
build = [
    "# FP32 — highest accuracy",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision fp32",
    "# FP16 — recommended for edge deployment",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision fp16",
    "# INT8 — fastest  (requires calibration images)",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision int8 --first true",
    "# DLA  — Xavier / Orin only",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision int8 --first true --dla 0",
]
bullet_box(sl, build, 0.5, 1.7, 12.33, 2.1, size=12, color=CODE_FG, bg=CODE_BG)

rect(sl, 0.5, 3.85, 12.33, 0.38, fill=RGBColor(0xFF, 0xF8, 0xE1), line=RGBColor(0xF0, 0xC0, 0x00))
txt(sl, "  Note:  --first true skips INT8 quantization on the first layer (accuracy-sensitive).",
    0.5, 3.88, 12.33, 0.32, size=12, color=RGBColor(0x7A, 0x55, 0x00))

section_label(sl, "Step 2 — Run Inference", 0.5, 4.35, 12.33, color=BLUE)
infer = [
    "# --d  Image directory   (Space: next image,  q: quit)",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision fp16 --d /path/to/images",
    "# --v  Video file  (MP4, AVI, ...)",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision fp16 --v /path/to/video.mp4",
    "# --cam  Live camera  (OpenCV device ID: 0 = /dev/video0)",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision fp16 --cam 0",
    "# Save results to disk",
    "./trt-lightnet --flagfile ../configs/CONFIG.txt --precision fp16 \\",
    "    --d images/ --save-detections true --save-detections-path ./output",
]
bullet_box(sl, infer, 0.5, 4.7, 12.33, 2.42, size=12, color=CODE_FG, bg=CODE_BG)


# ════════════════════════════════════════════════════════════════
# 7 — Configuration Files
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Configuration Files",
            "gflags-style plain text — one flag per line  ·  CLI flags override config")
page_footer(sl, 7)

section_label(sl, "Detection only", 0.5, 1.35, 5.9)
cfg1 = [
    "--onnx=/path/to/model.onnx",
    "--names=/path/to/classes.names",
    "--rgb=/path/to/colormap.colormap",
    "--precision=fp16",
    "--anchors=8,9,36,11,15,28,40,36,29,72",
    "--num_anchors=3",
    "--c=10",
    "--thresh=0.2",
    "--nms_thresh=0.6",
    "--cuda=true",
]
bullet_box(sl, cfg1, 0.5, 1.7, 5.9, 2.9, size=11, color=CODE_FG, bg=CODE_BG)

section_label(sl, "Two-stage (main + TLR subnet)", 6.9, 1.35, 6.0, color=BLUE)
cfg2 = [
    "# Main model",
    "--onnx=/path/to/main.onnx",
    "--names=main.names  --c=10  --thresh=0.2",
    "",
    "# Subnet (e.g. traffic light recognition)",
    "--subnet_onnx=/path/to/tlr.onnx",
    "--subnet_names=tlr.names",
    "--subnet_c=6  --subnet_thresh=0.2",
    "--target_names=trigger.names",
    "--batch_size=64",
]
bullet_box(sl, cfg2, 6.9, 1.7, 6.0, 2.9, size=11, color=CODE_FG, bg=CODE_BG)

section_label(sl, "Data File Formats", 0.5, 4.75, 12.4, color=GRAY_DARK)
table(sl,
      ["Extension", "Format", "Description"],
      [
          (".names",                 "One class name per line",       "Detection class labels"),
          (".colormap",              "R,G,B per line",                "Per-class visualization colors"),
          (".csv",                   "id, name, r, g, b, is_dynamic", "Segmentation class metadata"),
          ("calibration_images.txt", "One image path per line",       "INT8 calibration image list"),
      ],
      0.5, 5.1, 12.4, [2.2, 3.5, 6.7], row_size=11)


# ════════════════════════════════════════════════════════════════
# 8 — Multitask Inference
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Multitask Inference",
            "Detection · Segmentation · Depth — one forward pass")
page_footer(sl, 8)

# Pipeline
stages = [
    ("Input\nImage",    GRAY_MID),
    ("Backbone\n(LightNet)", GREEN),
    ("Neck\n(FPN)",      BLUE),
    ("Detection\nHead",  RGBColor(0xE0, 0x70, 0x00)),
    ("Seg\nHead",        RGBColor(0x80, 0x00, 0x90)),
    ("Depth\nHead",      RGBColor(0xC0, 0x00, 0x30)),
]
bw, bh, gap = 1.75, 0.65, 0.22
sx = 0.5
for i, (label, color) in enumerate(stages):
    lx = sx + i * (bw + gap)
    rect(sl, lx, 1.4, bw, bh, fill=color)
    txt(sl, label, lx, 1.4, bw, bh,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        txt(sl, "›", lx + bw, 1.52, gap + 0.08, 0.42,
            size=16, color=GRAY_MID, align=PP_ALIGN.CENTER)

# Output cards
outputs = [
    ("Object Detection",
     "Bounding boxes with class label,\nconfidence score, and optional keypoints.",
     RGBColor(0xE0, 0x70, 0x00)),
    ("Semantic Segmentation",
     "Per-pixel class labels rendered\nwith configurable colormap.",
     RGBColor(0x80, 0x00, 0x90)),
    ("Depth Estimation",
     "Dense depth map →\nBEV occupancy grid (requires fx, fy).",
     RGBColor(0xC0, 0x00, 0x30)),
]
for i, (title, desc, color) in enumerate(outputs):
    lx = 0.5 + i * 4.3
    rect(sl, lx, 2.3, 4.0, 0.38, fill=color)
    txt(sl, title, lx + 0.1, 2.32, 3.8, 0.34,
        size=13, bold=True, color=WHITE)
    rect(sl, lx, 2.68, 4.0, 0.88, fill=GRAY_LIGHT, line=GRAY_LINE)
    txt(sl, desc, lx + 0.12, 2.73, 3.76, 0.8, size=12, color=GRAY_DARK)

section_label(sl, "Config — detection + segmentation + depth", 0.5, 3.75, 12.33, color=GRAY_DARK)
cfg_mt = [
    "--onnx=multitask.onnx  --names=classes.names  --rgb=detect.colormap  --precision=fp16",
    "--anchors=8,9,36,11,15,28  --num_anchors=3  --c=8  --thresh=0.2  --cuda=true",
    "--mask=segmentation.csv              # segmentation",
    "--fx=1000.0  --fy=1000.0  --max_distance=80.0   # depth / BEV",
]
bullet_box(sl, cfg_mt, 0.5, 4.1, 12.33, 1.2, size=12, color=CODE_FG, bg=CODE_BG)

section_label(sl, "Python", 0.5, 5.45, 12.33, color=BLUE)
py_mt = [
    "lightnet.infer(frame, cuda=True)                        # run multitask inference",
    "bboxes = lightnet.get_bboxes()                          # detection results",
    "masks  = lightnet.get_masks_from_cpp()                  # segmentation masks",
]
bullet_box(sl, py_mt, 0.5, 5.8, 12.33, 0.9, size=12, color=CODE_FG, bg=CODE_BG)


# ════════════════════════════════════════════════════════════════
# 9 — Two-Stage Hierarchical Detection
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Two-Stage Hierarchical Detection",
            "Main model → crop ROIs → subnet classifier  (e.g. Traffic Light Recognition)")
page_footer(sl, 9)

flow = [
    ("Full Image", GRAY_MID),
    ("Main Model\nDetection",    GREEN),
    ("Crop ROIs\n(target class)", BLUE),
    ("Subnet\nBatch Infer",      RGBColor(0xE0, 0x70, 0x00)),
    ("Fine-grained\nResult",     RGBColor(0x80, 0x00, 0x90)),
]
bw2 = 2.1
sx2 = 0.5
for i, (label, color) in enumerate(flow):
    lx = sx2 + i * (bw2 + 0.28)
    rect(sl, lx, 1.4, bw2, 0.75, fill=color)
    txt(sl, label, lx, 1.4, bw2, 0.75,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        txt(sl, "›", lx + bw2, 1.52, 0.35, 0.52,
            size=18, color=GRAY_MID, align=PP_ALIGN.CENTER)

section_label(sl, "Main Model", 0.5, 2.4, 5.9)
main_items = [
    "• Detects all object categories",
    "• Instances matching --target_names trigger subnet",
    "  (e.g. traffic_light.names activates TLR subnet)",
    "",
    "Key flags:",
    "  --onnx           Main ONNX model path",
    "  --c              Number of main classes",
    "  --thresh         Detection threshold",
    "  --target_names   Trigger class list",
]
bullet_box(sl, main_items, 0.5, 2.75, 5.9, 3.6, size=12, color=GRAY_DARK, bg=GRAY_LIGHT)

section_label(sl, "Subnet Classifier", 6.9, 2.4, 6.0, color=BLUE)
sub_items = [
    "• Crops batch-processed at --batch_size (e.g. 64)",
    "• Higher resolution → finer classification",
    "• Outputs subclass per detected instance",
    "  e.g. red / amber / green / left / right arrow",
    "",
    "Key flags:",
    "  --subnet_onnx     Subnet ONNX model path",
    "  --subnet_c        Number of subnet classes",
    "  --subnet_thresh   Subnet confidence threshold",
    "  --batch_size      Crop batch size",
]
bullet_box(sl, sub_items, 6.9, 2.75, 6.0, 3.6, size=12, color=GRAY_DARK, bg=GRAY_LIGHT)


# ════════════════════════════════════════════════════════════════
# 10 — Uncertainty & BEV
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Uncertainty Estimation & BEV Projection")
page_footer(sl, 10)

section_label(sl, "Uncertainty Estimation  (--entropy)", 0.5, 1.35, 5.9)
unc = [
    "Computes Shannon entropy from softmax output:",
    "  H(x) = −Σ p(c|x) · log p(c|x)",
    "",
    "High entropy → low model confidence",
    "Useful for OOD / uncertainty-aware driving",
    "",
    "Output:",
    "  • Per-pixel entropy maps (heatmap)",
    "  • Per-class scalar entropy values",
    "",
    "Python API:",
    "  lightnet.make_entropy()",
    "  maps = lightnet.get_entropy_maps_from_cpp()",
    "  vals = lightnet.get_entropies()",
]
bullet_box(sl, unc, 0.5, 1.7, 5.9, 5.3, size=12, color=GRAY_DARK, bg=GRAY_LIGHT)

section_label(sl, "Bird's Eye View  (BEV) Projection", 6.9, 1.35, 6.0, color=BLUE)
bev = [
    "Back-projects depth map into top-down grid",
    "using pin-hole camera model:",
    "",
    "Pipeline:",
    "  1. Depth head → dense depth map",
    "  2. Back-projection with fx, fy",
    "  3. Height map → jet/magma colorized BEV",
    "  4. Optional road-plane smoothing (--smooth)",
    "",
    "Required config:",
    "  --fx            Horizontal focal length (px)",
    "  --fy            Vertical focal length (px)",
    "  --max_distance  Max depth range (m)",
    "",
    "Optional:",
    "  --smooth  Road segmentation-guided correction",
]
bullet_box(sl, bev, 6.9, 1.7, 6.0, 5.3, size=12, color=GRAY_DARK, bg=GRAY_LIGHT)


# ════════════════════════════════════════════════════════════════
# 11 — LiDAR Range Image Segmentation
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "LiDAR Range Image Segmentation",
            "Point cloud → range image → semantic segmentation")
page_footer(sl, 11)

stages3 = [
    ("LiDAR\nPoint Cloud\n(.bin)",           GRAY_MID),
    ("Range Image\nGeneration",               GREEN),
    ("TensorRT\nInference",                   BLUE),
    ("Segmentation\nMask",                    RGBColor(0xE0, 0x70, 0x00)),
    ("Entropy Map\n(optional)",               RGBColor(0x80, 0x00, 0x90)),
]
bw3, gap3 = 2.1, 0.27
sx3 = 0.5
for i, (label, color) in enumerate(stages3):
    lx = sx3 + i * (bw3 + gap3)
    rect(sl, lx, 1.4, bw3, 0.85, fill=color)
    txt(sl, label, lx, 1.4, bw3, 0.85,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages3) - 1:
        txt(sl, "›", lx + bw3, 1.56, gap3 + 0.1, 0.52,
            size=18, color=GRAY_MID, align=PP_ALIGN.CENTER)

section_label(sl, "CLI", 0.5, 2.45, 5.9, color=GRAY_DARK)
cli = [
    "./trt-lightnet \\",
    "  --flagfile configs/RangeImage.txt \\",
    "  --precision fp16 \\",
    "  --lidar \\",
    "  --camera_name CAM_FRONT_WIDE",
    "",
    "Key flags:",
    "  --lidar         Enable range image mode",
    "  --camera_name   Sensor identifier",
    "  --entropy       Enable uncertainty maps",
]
bullet_box(sl, cli, 0.5, 2.8, 5.9, 3.2, size=12, color=CODE_FG, bg=CODE_BG)

section_label(sl, "Python — T4 Dataset Demo", 6.9, 2.45, 6.0, color=BLUE)
py3 = [
    "python scripts/range_image_demo.py \\",
    "  --t4dataset /path/to/t4d \\",
    "  --camera-name CAM_FRONT_WIDE \\",
    "  --flagfile configs/RangeImage.txt \\",
    "  --save-segmentation \\",
    "  --save-uncertainty \\",
    "  --output-dir ./output",
    "",
    "Outputs:",
    "  {stem}_seg_{i}.png          segmentation mask",
    "  {stem}_uncertainty_{i}.png  entropy map",
    "  {stem}_range_image.png      raw range image",
]
bullet_box(sl, py3, 6.9, 2.8, 6.0, 3.2, size=12, color=CODE_FG, bg=CODE_BG)


# ════════════════════════════════════════════════════════════════
# 12 — Python API
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Python API — pylightnet",
            "ctypes-based wrapper  ·  cd python && pip install .")
page_footer(sl, 12)

section_label(sl, "Install & Demo Scripts", 0.5, 1.35, 4.5, color=GRAY_DARK)
inst = [
    "cd python",
    "pip install setuptools==68.2.2",
    "pip install .",
    "",
    "# Video demo",
    "python scripts/demo.py -f config.txt -v video.mp4",
    "",
    "# T4 dataset (LiDAR)",
    "python scripts/range_image_demo.py \\",
    "  -t4d /t4d -cam CAM_FRONT_WIDE -f cfg.txt",
]
bullet_box(sl, inst, 0.5, 1.7, 4.5, 3.1, size=11, color=CODE_FG, bg=CODE_BG)

section_label(sl, "API Usage", 5.3, 1.35, 7.6, color=BLUE)
api = [
    "import pylightnet",
    "",
    "config   = pylightnet.load_config('config.txt')",
    "lightnet = pylightnet.create_lightnet_from_config(config)",
    "",
    "lightnet.infer(frame, cuda=True)            # BGR numpy array",
    "",
    "# Detection",
    "bboxes = lightnet.get_bboxes()",
    "pylightnet.draw_bboxes_on_image(frame, bboxes, colormap, names)",
    "",
    "# Segmentation",
    "seg_data   = pylightnet.load_segmentation_data(config['mask'])",
    "lightnet.make_mask(lightnet.segmentation_to_argmax2bgr(seg_data))",
    "masks = lightnet.get_masks_from_cpp()",
    "",
    "# Uncertainty",
    "lightnet.make_entropy()",
    "entropy_maps = lightnet.get_entropy_maps_from_cpp()",
    "",
    "lightnet.destroy()                          # free GPU memory",
]
bullet_box(sl, api, 5.3, 1.7, 7.6, 5.3, size=11, color=CODE_FG, bg=CODE_BG)

section_label(sl, "Docker Test", 0.5, 4.9, 4.5, color=GRAY_DARK)
bullet_box(sl, ["make test-pylightnet"],
           0.5, 5.25, 4.5, 0.45, size=11, color=CODE_FG, bg=CODE_BG)


# ════════════════════════════════════════════════════════════════
# 13 — Performance Optimization
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Performance Optimization")
page_footer(sl, 13)

opts = [
    ("FP16",              "--precision fp16",
     "~2× faster than FP32.\nNegligible accuracy loss.\nRecommended for most deployments."),
    ("INT8",              "--precision int8 --first true\n--calib Entropy",
     "~4× faster than FP32.\nRequires 50–500 calibration images.\n--first true protects first layer."),
    ("2:4 Sparsity",      "--sparse true",
     "50% weight reduction.\nRequires Ampere+ GPU.\nModel must be sparsity-aware trained."),
    ("NVDLA",             "--dla 0",
     "Offload to DLA accelerator.\nXavier / Orin only.\nReduces GPU power consumption."),
    ("Per-layer Profile", "--profile",
     "Prints per-layer latency.\nHelps identify bottleneck layers\nfor precision tuning."),
    ("GPU Preprocess",    "infer(frame, cuda=True)",
     "Resize + normalize on GPU.\nAvoids CPU↔GPU transfer.\nEnabled via CUDA kernels."),
]
for i, (title, flag, desc) in enumerate(opts):
    col, row = i % 3, i // 3
    lx = 0.4 + col * 4.3
    ty = 1.4 + row * 2.85
    rect(sl, lx, ty, 4.05, 2.65, fill=WHITE, line=GRAY_LINE)
    rect(sl, lx, ty, 4.05, 0.38, fill=GREEN if row == 0 else BLUE)
    txt(sl, title, lx + 0.1, ty + 0.04, 3.85, 0.3,
        size=13, bold=True, color=WHITE)
    bullet_box(sl, [flag], lx, ty + 0.38, 4.05, 0.6,
               size=11, color=CODE_FG, bg=CODE_BG)
    txt(sl, desc, lx + 0.12, ty + 1.08, 3.81, 1.48,
        size=12, color=GRAY_DARK)


# ════════════════════════════════════════════════════════════════
# 14 — Command-Line Reference
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Command-Line Reference",
            "Full list in  src/config_parser.cpp")
page_footer(sl, 14)

cols3 = [
    ("Core Flags", GREEN, [
        ("--flagfile",   "string", "Path to config file  (required)"),
        ("--precision",  "string", "fp32 / fp16 / int8"),
        ("--onnx",       "string", "ONNX model path"),
        ("--names",      "string", "Class names file"),
        ("--c",          "int",    "Number of detection classes"),
        ("--thresh",     "float",  "Detection confidence threshold"),
        ("--nms_thresh", "float",  "NMS IoU threshold"),
    ]),
    ("Input Source  (choose one)", BLUE, [
        ("--d",          "string", "Image directory (Space: next,  q: quit)"),
        ("--v",          "string", "Video file  (MP4, AVI, ...)"),
        ("--cam",        "int",    "Live camera device ID  (0 = /dev/video0)"),
        ("--save-detections",      "bool",   "Save output images"),
        ("--save-detections-path", "string", "Output directory"),
        ("--profile",    "bool",   "Print per-layer latency"),
    ]),
    ("Model Optimization", RGBColor(0xE0, 0x70, 0x00), [
        ("--first",  "bool",   "Skip INT8 quant on first layer  (recommended: true)"),
        ("--last",   "bool",   "Skip INT8 quant on last layer"),
        ("--sparse", "bool",   "2:4 structured sparsity  (Ampere+ only)"),
        ("--dla",    "int",    "DLA core 0 or 1  (Xavier / Orin only)"),
        ("--calib",  "string", "INT8 calib: Entropy / MinMax / Percentile"),
        ("--calibration_images", "string", "Calibration image list"),
    ]),
]
col_xs = [0.35, 4.65, 8.95]
for ci, (sec, color, rows) in enumerate(cols3):
    lx = col_xs[ci]
    cw = 4.1
    rect(sl, lx, 1.35, cw, 0.35, fill=color)
    txt(sl, sec, lx + 0.1, 1.37, cw - 0.15, 0.31,
        size=12, bold=True, color=WHITE)
    for ri, (flag, ftype, desc) in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else GRAY_LIGHT
        ty = 1.7 + ri * 0.72
        rect(sl, lx, ty, cw, 0.72, fill=bg, line=GRAY_LINE, lw=Pt(0.5))
        txt(sl, flag, lx + 0.1, ty + 0.06, cw - 0.15, 0.28,
            size=12, bold=True, color=color)
        txt(sl, f"[{ftype}]  {desc}", lx + 0.1, ty + 0.36, cw - 0.15, 0.28,
            size=11, color=GRAY_MID)


# ════════════════════════════════════════════════════════════════
# 15 — Tips: Model Optimization
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Tips — Model Optimization Flags",
            "--first  ·  --sparse  ·  --dla")
page_footer(sl, 15, total=16)

# ── --first ──────────────────────────────────────────────────────
section_label(sl, "--first  /  --last", 0.5, 1.35, 4.0)
first_items = [
    "The first (and last) conv layers process raw",
    "pixel values and have a much wider activation",
    "dynamic range than intermediate layers.",
    "",
    "Quantizing them to INT8 often causes a",
    "noticeable accuracy drop.",
    "",
    "→  --first true keeps the first layer in",
    "   FP32/FP16 while all others run in INT8.",
    "",
    "Recommended for all INT8 deployments:",
    "  --precision int8 --first true",
]
bullet_box(sl, first_items, 0.5, 1.7, 4.0, 5.1,
           size=12, color=GRAY_DARK, bg=GRAY_LIGHT)

# ── --sparse ─────────────────────────────────────────────────────
section_label(sl, "--sparse  (2:4 Structured Sparsity)", 4.9, 1.35, 4.0, color=BLUE)

# mini diagram
for i, (val, zero) in enumerate([("0.3", False), ("-0.7", True),
                                   ("0.1", False), ("-0.5", True)]):
    lx = 4.9 + i * 0.92
    c = GRAY_LIGHT if zero else WHITE
    rect(sl, lx, 1.78, 0.88, 0.42, fill=c, line=GRAY_LINE)
    txt(sl, "0.0" if zero else val, lx, 1.78, 0.88, 0.42,
        size=12, bold=zero, color=GRAY_MID if zero else BLACK,
        align=PP_ALIGN.CENTER)
txt(sl, "2 zeros forced per 4 weights", 4.9, 2.25, 3.7, 0.3,
    size=10, color=GRAY_MID, italic=True)

sparse_items = [
    "NVIDIA Ampere sparse tensor cores execute",
    "pruned weights at ~2× GPU throughput.",
    "",
    "Requirements:",
    "  • Ampere GPU or later  (RTX 30xx / A-series)",
    "  • Model trained with sparsity-aware",
    "    training  (NVIDIA ASP / PyTorch)",
    "",
    "Enabling --sparse on a dense model will",
    "degrade accuracy — train first.",
    "",
    "  --precision fp16 --sparse true",
]
bullet_box(sl, sparse_items, 4.9, 2.6, 4.0, 4.2,
           size=12, color=GRAY_DARK, bg=GRAY_LIGHT)

# ── --dla ────────────────────────────────────────────────────────
section_label(sl, "--dla  (NVDLA — Xavier / Orin only)", 9.3, 1.35, 3.7, color=RGBColor(0xE0, 0x70, 0x00))
dla_items = [
    "DLA is a dedicated fixed-function HW",
    "accelerator inside Jetson Xavier / Orin.",
    "",
    "Benefits:",
    "  • Frees GPU for other workloads",
    "  • Significantly lower power draw",
    "  • Two independent cores (0 and 1)",
    "",
    "Constraints:",
    "  • INT8 precision only",
    "  • Unsupported layers fall back to GPU",
    "  • NOT available on desktop/server GPUs",
    "",
    "  --precision int8 --first true --dla 0",
    "",
    "Ref: TIER IV Tech Blog (docswell 2023)",
    "  docswell.com/s/TIER_IV/",
    "  KGX2L8-2023-07-24-120048#p27",
]
bullet_box(sl, dla_items, 9.3, 1.7, 3.7, 5.1,
           size=11, color=GRAY_DARK, bg=GRAY_LIGHT)


# ════════════════════════════════════════════════════════════════
# 16 — Project Structure & References
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
page_header(sl, "Project Structure & References")
page_footer(sl, 16, total=16)

section_label(sl, "Directory Structure", 0.5, 1.35, 6.5, color=GRAY_DARK)
struct = [
    "trt-lightnet/",
    "├── configs/    config files (.txt) & precision tables (.json)",
    "├── data/       colormaps, class names, segmentation CSVs",
    "├── include/",
    "│   ├── tensorrt_lightnet/  inference engine headers",
    "│   ├── tensorrt_common/    TensorRT wrapper",
    "│   ├── cuda_utils/         CUDA memory helpers",
    "│   ├── sensor/             calibration parsing",
    "│   └── pcdUtils/           point cloud utilities",
    "├── src/",
    "│   ├── lightnet_detector.cpp   main executable",
    "│   ├── tensorrt_lightnet/      inference implementation",
    "│   ├── preprocess.cu           GPU preprocessing kernels",
    "│   └── config_parser.cpp       gflags config parsing",
    "├── python/",
    "│   ├── _pylightnet.py          ctypes wrapper",
    "│   └── scripts/                demo & utility scripts",
    "└── Dockerfile_x86 / _aarch64",
]
bullet_box(sl, struct, 0.5, 1.7, 6.5, 5.1, size=10, color=CODE_FG, bg=CODE_BG)

section_label(sl, "References", 7.4, 1.35, 5.5, color=BLUE)
refs = [
    "[1]  LightNet — base CNN architecture",
    "     github.com/daniel89710/lightNet",
    "",
    "[2]  TensorRT — NVIDIA inference optimizer",
    "     developer.nvidia.com/tensorrt",
    "",
    "[3]  2:4 Structured Sparsity",
    "     NVIDIA Developer Blog (2022)",
    "",
    "[4]  NVDLA — nvdla.org",
    "",
    "[5]  INT8 Quantization Aware Training",
    "     NVIDIA Developer Blog (2021)",
    "",
    "[6]  lightNet-TR (original)",
    "     github.com/daniel89710/trt-lightnet",
    "",
    "Repository:",
    "  github.com/tier4/trt-lightnet",
    "",
    "License: Apache 2.0",
]
bullet_box(sl, refs, 7.4, 1.7, 5.5, 5.1, size=11, color=GRAY_DARK, bg=GRAY_LIGHT)


# ── Save ─────────────────────────────────────────────────────────
out = "/home/danumeda/work/trt-lightnet/docs/trt-lightnet-overview.pptx"
prs.save(out)
print(f"Saved: {out}")
